"""
Parse PPT / PDF / TXT files into plain text for LLM consumption.
Uses PyMuPDF for PDF (better quality than PyPDF2) with text cleaning.
"""

import os
import re


def parse_file(path: str) -> str:
    """Return extracted plain text from the given file."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pptx", ".ppt"):
        return _parse_pptx(path)
    elif ext == ".pdf":
        return _parse_pdf(path)
    elif ext in (".txt", ".md"):
        return _parse_txt(path)
    else:
        raise ValueError(f"不支持的文件格式：{ext}，请上传 PPT/PDF/TXT 文件。")


# Alias for member2-compatible call sites
def extract_text(path: str) -> str:
    return parse_file(path)


def _clean_text(text: str) -> str:
    """Remove page numbers, TOC dotted lines, and blank lines."""
    page_num_pat = re.compile(r'^\s*\d+\s*$')
    toc_pat = re.compile(r'\.{4,}|…{2,}')
    cleaned = []
    for line in text.split('\n'):
        s = line.strip()
        if not s:
            continue
        if page_num_pat.match(s):
            continue
        if s in ("目录", "Contents") or toc_pat.search(s):
            continue
        cleaned.append(s)
    return '\n'.join(cleaned)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for shape in (shape for slide in prs.slides for shape in slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            lines.append(shape.text.strip())
    return _clean_text('\n'.join(lines))


def _parse_pdf(path: str) -> str:
    """Use PyMuPDF for better text extraction quality."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        pages = [page.get_text("text") for page in doc]
        return _clean_text('\n'.join(pages))
    except ImportError:
        pass
    # Fallback to PyPDF2 if PyMuPDF not installed
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = [page.extract_text() or "" for page in reader.pages]
        return _clean_text('\n'.join(pages))
    except Exception as e:
        raise RuntimeError(f"PDF解析失败：{e}")


def _parse_txt(path: str) -> str:
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise RuntimeError("无法读取文本文件，请确认编码格式。")


def truncate_text(text: str, max_chars: int = 12000) -> str:
    """Truncate to avoid exceeding LLM context limits."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[内容已截断，原文共 {len(text)} 字符]"