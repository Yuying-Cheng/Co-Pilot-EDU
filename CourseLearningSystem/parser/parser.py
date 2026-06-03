import re
import os
from pptx import Presentation
import fitz  # PyMuPDF

try:
    from docx import Document
except ImportError:
    Document = None


def clean_text(text: str) -> str:
    """
    文本清洗模块
    功能：去页码、去目录、去空行
    """
    lines = text.split('\n')
    cleaned_lines = []

    page_num_pattern = re.compile(r'^\s*\d+\s*$')
    toc_pattern = re.compile(r'\.{4,}|…{2,}')

    for line in lines:
        stripped_line = line.strip()

        if not stripped_line:
            continue
        if page_num_pattern.match(stripped_line):
            continue
        if (stripped_line == "目录") or ("Contents" in stripped_line) or toc_pattern.search(stripped_line):
            continue

        cleaned_lines.append(stripped_line)

    return '\n'.join(cleaned_lines)


def read_ppt(file_path: str) -> str:
    """PPT解析模块"""
    try:
        prs = Presentation(file_path)
        full_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)

        raw_text = '\n'.join(full_text)
        return clean_text(raw_text)

    except Exception as e:
        print(f"❌ PPT 解析失败: {e}")
        return ""


def read_pdf(file_path: str) -> str:
    """PDF解析模块"""
    try:
        doc = fitz.open(file_path)
        full_text = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            full_text.append(text)

        raw_text = '\n'.join(full_text)
        return clean_text(raw_text)

    except Exception as e:
        print(f"❌ PDF 解析失败: {e}")
        return ""


def read_text_file(file_path: str) -> str:
    """TXT/Markdown 讲义解析模块"""
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return clean_text(f.read())
        except UnicodeDecodeError:
            continue
        except Exception as e:
            print(f"❌ 文本文件解析失败: {e}")
            return ""
    print("❌ 文本文件解析失败: 无法识别文件编码")
    return ""


def read_docx(file_path: str) -> str:
    """DOCX 讲义解析模块"""
    if Document is None:
        print("❌ DOCX 解析失败: 缺少 python-docx 依赖")
        return ""
    try:
        doc = Document(file_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return clean_text("\n".join(parts))
    except Exception as e:
        print(f"❌ DOCX 解析失败: {e}")
        return ""


def extract_text(file_path: str) -> str:
    """统一的文件解析入口（根据后缀自动分发）"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in ['.ppt', '.pptx']:
        return read_ppt(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.txt', '.md']:
        return read_text_file(file_path)
    elif ext == '.docx':
        return read_docx(file_path)
    else:
        raise ValueError(f"❌ 不支持的文件格式: {ext}。请上传 PPT、PDF、DOCX、TXT 或 MD。")
