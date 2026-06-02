"""
Parse PPT / PDF / TXT files into plain text for LLM consumption.
"""

import os
import re

def parse_file(path: str) -> str:
    """Return extracted plain text from the given file."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pptx", ".ppt"):
        return _parse_pptx(path)
    elif ext == ".pdf":
        return _parse_pdf(path)
    elif ext in (".txt", ".md"):
        return _parse_txt(path)
    else:
        raise ValueError(f"不支持的文件格式：{ext}，请上传 PPT/PDF/TXT 文件。")


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            lines.append(f"【第{i}页】")
            lines.extend(slide_texts)
    return "\n".join(lines)


def _parse_pdf(path: str) -> str:
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"【第{i}页】\n{text.strip()}")
            return "\n\n".join(pages)
    except Exception as e:
        raise RuntimeError(f"PDF解析失败：{e}")


def _parse_txt(path: str) -> str:
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise RuntimeError("无法读取文本文件，请确认编码格式。")


def truncate_text(text: str, max_chars: int = 12000) -> str:
    """Truncate to avoid exceeding LLM context limits."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[内容已截断，原文共 {len(text)} 字符]"
